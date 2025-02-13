import os
import docx
import pptx
import shutil
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer
from modules.extras import hashx

def get_file_content (path, current_dir, collection):
    # get the file extension
    extension = path.split('.')[-1]
    # make it lowercase
    extension = extension.lower()
    # copy the file to assets
    new_path = os.environ["COLLECTIONS_PATH"] + f"{collection}/assets/{hashx(path.split('/')[-1])}.{extension}"
    # the content
    content = []
    # a flag
    accepted = True

######### if pdf #############################################
    
    if extension == "pdf":
        try:
            # for each page in pdf
            for page_layout in extract_pages(path):
                # save every page
                page = ""
                # for each element in page
                for element in page_layout:
                    # if is text
                    if isinstance(element, LTTextContainer):
                        # then save it on content
                        page += element.get_text()
                # save the page
                content.append(page)
        # in case of error
        except Exception as e:
            print(f"Error Reading PDF: {e}")

######### if word ############################################

    elif extension == "docx":
        try:
            page = ""
            # try to open it
            document = docx.Document(path)
            # for each paragraph
            for i, paragraph in enumerate(document.paragraphs):
                # save the paragraph
                page += paragraph.text + "\n\n"
                # and every three pages save it
                if (i+1) % 3 == 0:
                    content.append(page)
                    page = ""
        # in case of error
        except Exception as e:
            print(f"Error Reading DOCX: {e}")

######### if power point #####################################

    elif extension == "pptx":
        try:
            # try to open the presentation
            presentation = pptx.Presentation(path)
            # for each slide
            for slide in presentation.slides:
                # set the text for every slide
                text = ""
                # for each item on slide
                for shape in slide.shapes:
                    # if it has text
                    if shape.has_text_frame:
                        # sum it to text
                        text += shape.text + ". "
                # and append to content
                content.append(text.strip())
        except Exception as e:
            print(f"Error Reading PPTX: {e}")

######### if txt ############################################

    elif extension in ["txt", "md"]:
        try:
            text = ""
            # open the txt file
            with open(path, "r", encoding="utf-8") as file:
                # for each line
                for line in file:
                    # add the line to text
                    text += line
            # save text in content
            content.append(text)
        except Exception as e:
            print(f"Error Reading TXT: {e}")

##############################################################

    else:
        print(f"File not supported, only .pdf .docx .pptx .txt .md")
        accepted = False

    if accepted:
        print(f"processing {path} with {len(content)} items")


    return content, new_path